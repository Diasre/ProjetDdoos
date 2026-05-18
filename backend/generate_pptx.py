import os
import subprocess
import sys

print("Installation de la librairie python-pptx si nécessaire...")
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

print("Génération du PowerPoint en cours...")

prs = Presentation()

# Theming (Darkish)
def apply_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate-900

def set_text(shape, text, size=18, bold=False, rgb=(255, 255, 255)):
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*rgb)

# Slide 1 : Titre
slide = prs.slides.add_slide(prs.slide_layouts[0])
apply_theme(slide)
set_text(slide.shapes.title, "IPS Intelligent Anti-DDoS", 40, True)
set_text(slide.placeholders[1], "Couplage d'Intelligence Artificielle LSTM et de Prévention Cisco\n\nProjet de Fin d'Études / UIT 2026", 20, False, (200, 200, 200))

# Helper pour slides classiques
def add_bullet_slide(title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    
    # Titre
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248) # Sky blue
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Contenu
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = points[0]
    tf.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    tf.paragraphs[0].font.size = Pt(20)
    
    for pt in points[1:]:
        p = tf.add_paragraph()
        p.text = pt
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.font.size = Pt(20)

# Slide 2
add_bullet_slide(
    "1. Contexte & Problématique",
    [
        "• Évolution redoutable : amplification par réflexion (DrDoS).",
        "• Les pare-feux classiques statiques sont aveugles face aux flux massifs complexes.",
        "• Besoin Vital : Systèmes adaptatifs, autonomes et d'isolation instantanée."
    ]
)

# Slide 3
add_bullet_slide(
    "2. Architecture Globale à Trois Couches",
    [
        "• L'Intelligence (Backend API) : FastAPI et TensorFlow (LSTM) ultra-rapide.",
        "• L'Action (Cisco IPS) : Scripts Netmiko modifiant l'ACL étendue IA_BLOCK_DDOS.",
        "• L'Observation (React UI) : Tableau de bord cartographique avec IA intégrée."
    ]
)

# Slide 4
add_bullet_slide(
    "3. Le Cerveau IA - Le modèle LSTM",
    [
        "• Base d'Entraînement : CIC-DDoS2019 (référence académique de pointe).",
        "• Algorithme LSTM : Détecte les anomalies dans la temporalité du trafic réseau.",
        "• Inférence Temps Réel : Extraction par Scapy et analyse en < 1.2ms."
    ]
)

# Slide 5
add_bullet_slide(
    "4. Bouclier Heuristique DrDoS Actif",
    [
        "Couverture des amplifications majeures (Taux de confiance 99%) :",
        "• LDAP (Port 389), MSSQL (Port 1434), NetBIOS (Port 137)",
        "• Portmap (Port 111), SSDP (Port 1900), NTP (Port 123)",
        "• Saturation Brute : UDP Flood, SYN-Flood, Brute Force."
    ]
)

# Slide 6
add_bullet_slide(
    "5. Action Cisco Zéro-Délai",
    [
        "• Règle de blocage autonome (Zero-Touch).",
        "• Extraction de l'IP du pirate depuis le flux empoisonné.",
        "• Connexion SSH locale (192.168.56.103) et écriture de la commande :",
        "  'deny ip host [IP] any log'",
        "• Filtrage préventif des réseaux locaux pour sécuriser l'administration."
    ]
)

# Slide 7
add_bullet_slide(
    "6. Centre de Commandement React",
    [
        "• Dashboard Modern & Dark Premium.",
        "• Géolocalisation dynamique des menaces via l'API (ip-api).",
        "• Rétrospective Temporelle : Base de données SQLite3 pour le Journal Annuel.",
        "• Génération de rapports PDF certifiant les audits."
    ]
)

# Slide 8
add_bullet_slide(
    "7. Démonstration & Simulation",
    [
        "• Construction de robots de crash-test locaux (Bots Multi-Threads).",
        "• Injection API directe contournant le pare-feu standard Windows.",
        "• Détection algorithmique en direct sur le Dashboard.",
        "• Disjonction physique du port du routeur Cisco."
    ]
)

# Slide 9
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
apply_theme(slide9)
slide9.shapes.title.text = "Merci de votre attention"
slide9.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(16, 185, 129) # Emerald
tf = slide9.placeholders[1].text_frame
tf.text = "Place à la Démonstration Technique Live."
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Presentation_Jury.pptx")
prs.save(out_path)
print(f"\n[SUCCES] Présentation PowerPoint générée dans : {out_path}")
